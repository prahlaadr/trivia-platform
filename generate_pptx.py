"""Generate a PPTX presentation from a parsed Quiz."""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from parser import parse_docx, Quiz, Round


# Colors
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)    # dark navy
BG_ROUND = RGBColor(0x16, 0x21, 0x3E)   # slightly lighter navy
ACCENT = RGBColor(0xE8, 0x4D, 0x5A)     # coral red
ACCENT2 = RGBColor(0x4E, 0xC9, 0xB0)    # teal green
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)


def set_slide_bg(slide, color: RGBColor):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_box(slide, left, top, width, height, lines: list[tuple[str, int, RGBColor, bool]],
                      alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple formatted lines. Each line is (text, font_size, color, bold)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox


def make_title_slide(prs: Presentation, quiz: Quiz):
    """Create the title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                 f"Pub Quiz #{quiz.quiz_number}", font_size=48, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3), Inches(8), Inches(0.8),
                 quiz.date, font_size=28, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Round overview
    lines = []
    for r in quiz.rounds:
        joker_tag = " (Joker)" if r.joker_eligible else ""
        lines.append((f"Round {r.number} — {r.title}{joker_tag}", 18, WHITE, False))
    add_multiline_box(slide, Inches(2.5), Inches(4), Inches(5), Inches(3), lines,
                      alignment=PP_ALIGN.LEFT)


def make_round_title_slide(prs: Presentation, rnd: Round):
    """Create a round title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_ROUND)

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.5),
                 f"ROUND {rnd.number}", font_size=20, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.5),
                 rnd.title, font_size=44, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Scoring info
    pts_text = f"{rnd.points_per_question} Point{'s' if rnd.points_per_question > 1 else ''} Per Correct Answer"
    if rnd.round_type == "progressive":
        pts_text = "10 → 8 → 6 → 4 → 2 Points"
    joker_text = "Joker Is IN PLAY" if rnd.joker_eligible else "Joker Is NOT In Play"
    joker_color = ACCENT2 if rnd.joker_eligible else ACCENT

    add_text_box(slide, Inches(1), Inches(3.5), Inches(8), Inches(0.5),
                 pts_text, font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.5),
                 joker_text, font_size=20, color=joker_color, bold=True,
                 alignment=PP_ALIGN.CENTER)

    if rnd.theme_description:
        add_text_box(slide, Inches(1), Inches(5), Inches(8), Inches(1.5),
                     rnd.theme_description, font_size=16, color=GRAY,
                     alignment=PP_ALIGN.CENTER)


def make_questions_slide(prs: Presentation, rnd: Round):
    """Create a slide showing all questions for the round."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(9), Inches(0.4),
                 f"Round {rnd.number} — {rnd.title}", font_size=16, color=ACCENT, bold=True)

    # Build question lines
    questions = [q for q in rnd.questions if not q.is_internet_only]
    lines = []
    for q in questions:
        lines.append((f"{q.number}. {q.text}", 14, WHITE, False))
        if q.choices:
            choices_text = "    " + "  |  ".join(
                c.split(". ", 1)[1] if ". " in c else c for c in q.choices
            )
            # Remove correct markers for question slide
            choices_text = choices_text.replace("✅", "").replace("(Correct)", "").replace("(✅Correct)", "")
            lines.append((choices_text, 12, GRAY, False))

    if not lines:
        return  # video rounds with no displayable questions

    add_multiline_box(slide, Inches(0.3), Inches(0.6), Inches(9.4), Inches(6.5), lines)


def make_qa_slides(prs: Presentation, rnd: Round):
    """Create individual question → answer reveal slides."""
    questions = [q for q in rnd.questions if not q.is_internet_only]

    for q in questions:
        # Question slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, BG_DARK)

        add_text_box(slide, Inches(0.3), Inches(0.15), Inches(9), Inches(0.4),
                     f"Round {rnd.number} — {rnd.title}", font_size=14, color=GRAY)
        add_text_box(slide, Inches(0.8), Inches(0.7), Inches(1), Inches(1),
                     str(q.number), font_size=72, color=ACCENT, bold=True)
        add_text_box(slide, Inches(0.5), Inches(2), Inches(9), Inches(3),
                     q.text, font_size=28, color=WHITE,
                     alignment=PP_ALIGN.CENTER)

        if q.choices:
            choice_lines = [(c, 22, GRAY, False) for c in q.choices]
            add_multiline_box(slide, Inches(2), Inches(4.5), Inches(6), Inches(2.5),
                              choice_lines, alignment=PP_ALIGN.LEFT)

        # Answer slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, BG_ROUND)

        add_text_box(slide, Inches(0.3), Inches(0.15), Inches(9), Inches(0.4),
                     f"Round {rnd.number} — {rnd.title}", font_size=14, color=GRAY)
        add_text_box(slide, Inches(0.8), Inches(0.7), Inches(1), Inches(1),
                     str(q.number), font_size=72, color=ACCENT, bold=True)
        add_text_box(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(2),
                     q.text, font_size=22, color=GRAY,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(2),
                     q.answer, font_size=40, color=ACCENT2, bold=True,
                     alignment=PP_ALIGN.CENTER)


def make_progressive_slides(prs: Presentation, rnd: Round):
    """Create slides for progressive reveal rounds (Guess Who/Where)."""
    # One slide per clue, accumulating
    for i, clue in enumerate(rnd.clues):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, BG_DARK)

        add_text_box(slide, Inches(0.3), Inches(0.15), Inches(9), Inches(0.4),
                     f"Round {rnd.number} — {rnd.title}", font_size=14, color=GRAY)

        # Points badge
        add_text_box(slide, Inches(3.5), Inches(0.6), Inches(3), Inches(0.8),
                     f"{clue['points']} POINTS", font_size=32, color=YELLOW, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Show all clues revealed so far
        lines = []
        for j in range(i + 1):
            c = rnd.clues[j]
            color = WHITE if j == i else GRAY  # current clue is bright
            lines.append((f"[{c['points']}pts] {c['text']}", 16, color, j == i))

        add_multiline_box(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(5), lines)

    # Answer reveal
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_ROUND)

    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(9), Inches(0.4),
                 f"Round {rnd.number} — {rnd.title}", font_size=14, color=GRAY)
    add_text_box(slide, Inches(1), Inches(2), Inches(8), Inches(1),
                 "ANSWER", font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3), Inches(8), Inches(2),
                 rnd.progressive_answer, font_size=48, color=ACCENT2, bold=True,
                 alignment=PP_ALIGN.CENTER)


def make_video_round_slides(prs: Presentation, rnd: Round):
    """Create slides for video rounds — just the answers list for reveal."""
    # Answers reveal slide
    if rnd.questions:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, BG_ROUND)

        add_text_box(slide, Inches(0.3), Inches(0.15), Inches(9), Inches(0.4),
                     f"Round {rnd.number} — {rnd.title} — ANSWERS", font_size=16, color=ACCENT, bold=True)

        # Video round: the "text" field IS the answer (no separate Q/A)
        lines = [(f"{q.number}. {q.answer or q.text}", 20, ACCENT2, False) for q in rnd.questions]
        add_multiline_box(slide, Inches(1), Inches(0.8), Inches(8), Inches(6), lines)


def make_tie_breaker_slides(prs: Presentation, quiz: Quiz):
    """Create tie breaker slides."""
    if not quiz.tie_breaker_question:
        return

    # Question
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, Inches(1), Inches(1), Inches(8), Inches(1),
                 "TIE BREAKER", font_size=36, color=YELLOW, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(2.5), Inches(9), Inches(3),
                 quiz.tie_breaker_question, font_size=28, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # Answer
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_ROUND)
    add_text_box(slide, Inches(1), Inches(1), Inches(8), Inches(1),
                 "TIE BREAKER", font_size=36, color=YELLOW, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(2.5), Inches(9), Inches(2),
                 quiz.tie_breaker_question, font_size=22, color=GRAY,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(4.5), Inches(9), Inches(2),
                 quiz.tie_breaker_answer, font_size=40, color=ACCENT2, bold=True,
                 alignment=PP_ALIGN.CENTER)


def generate_pptx(quiz: Quiz, output_path: str | Path) -> Path:
    """Generate a PPTX file from a Quiz object."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    make_title_slide(prs, quiz)

    # Each round
    for rnd in quiz.rounds:
        make_round_title_slide(prs, rnd)

        if rnd.round_type == "progressive":
            make_progressive_slides(prs, rnd)
        elif rnd.round_type == "video":
            make_video_round_slides(prs, rnd)
        else:
            # Questions overview slide
            make_questions_slide(prs, rnd)
            # Individual Q → A slides
            make_qa_slides(prs, rnd)

    # Tie breaker
    make_tie_breaker_slides(prs, quiz)

    output = Path(output_path)
    prs.save(str(output))
    return output


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python generate_pptx.py <docx_file_or_folder> [output_dir]")
        sys.exit(1)

    path = Path(sys.argv[1])
    output_dir = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("output")
    output_dir.mkdir(exist_ok=True)

    if path.is_dir():
        from parser import parse_folder
        quizzes = parse_folder(path)
        for quiz in quizzes:
            out = output_dir / f"Quiz_{quiz.quiz_number}.pptx"
            generate_pptx(quiz, out)
            print(f"Generated: {out}")
    else:
        quiz = parse_docx(path)
        out = output_dir / f"Quiz_{quiz.quiz_number}.pptx"
        generate_pptx(quiz, out)
        print(f"Generated: {out}")
